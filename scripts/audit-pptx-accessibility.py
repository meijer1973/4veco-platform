#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
audit-pptx-accessibility.py — Audit a PPTX for L1.5D v2 B8 accessibility.

Usage: python scripts/audit-pptx-accessibility.py <pptx-path>

Reports:
1. Every text run with explicit font.size, and counts of runs below 14pt
   and 18pt thresholds.
2. Contrast pairs the deck would need to satisfy WCAG AA (4.5:1 body,
   3.0:1 large). The audit reads only the explicit run-level
   fontSize+color where python-pptx exposes them. Inherited values
   from masters/layouts are not introspected by this script (would
   require XML walk).
3. Per-text-run color value + the slide-master background color
   recorded for the slide, so the operator can eyeball borderline
   pairs.

Exit code: 0 if no violations; 1 if any run below 14pt or any explicit
fg/bg pair below 3.0:1.
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation


# WCAG relative-luminance and contrast --------------------------------

def _channel(c):
    c = c / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def luminance(rgb_hex):
    """Take a 'RRGGBB' string, return WCAG relative luminance."""
    h = rgb_hex.lstrip('#')
    r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
    return 0.2126 * _channel(r) + 0.7152 * _channel(g) + 0.0722 * _channel(b)


def contrast(fg_hex, bg_hex):
    a = luminance(fg_hex); b = luminance(bg_hex)
    light = max(a, b); dark = min(a, b)
    return (light + 0.05) / (dark + 0.05)


# Master background lookup --------------------------------------------

def slide_bg_hex(slide):
    """Best-effort: read the slide's master/layout background color.

    pptxgenjs defines masters with `background: { color: PC.indigoDeep }`
    or similar, which lands as an explicit `<p:bg>` solidFill on the
    slide-master. Returns the hex string (without #) or None.
    """
    try:
        # Walk slide -> slide_layout -> slide_master, return first explicit bg.
        for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
            bg = src.background
            fill = bg.fill
            if fill.type is not None and hasattr(fill, 'fore_color'):
                fc = fill.fore_color
                if fc.type is not None and hasattr(fc, 'rgb') and fc.rgb is not None:
                    return str(fc.rgb)
    except Exception:
        pass
    return None


# Run walker ----------------------------------------------------------

def walk_text_runs(slide):
    """Yield (shape_name, paragraph_idx, run_idx, run, text)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for p_idx, p in enumerate(tf.paragraphs):
            for r_idx, r in enumerate(p.runs):
                yield (shape.name or '?', p_idx, r_idx, r, r.text or '')


def run_size_pt(run):
    try:
        sz = run.font.size
        if sz is None:
            return None
        return sz.pt if hasattr(sz, 'pt') else int(sz) / 12700.0
    except Exception:
        return None


def run_color_hex(run):
    try:
        c = run.font.color
        if c is None: return None
        if c.type is None: return None
        if hasattr(c, 'rgb') and c.rgb is not None:
            return str(c.rgb)
    except Exception:
        pass
    return None


# Main ----------------------------------------------------------------

def audit(pptx_path):
    pres = Presentation(pptx_path)

    below_14 = []
    below_18 = []
    contrast_violations = []
    total_runs = 0

    for s_idx, slide in enumerate(pres.slides, 1):
        bg = slide_bg_hex(slide)
        for (sh, p_idx, r_idx, run, text) in walk_text_runs(slide):
            if not text.strip():
                continue
            total_runs += 1
            sz = run_size_pt(run)
            fg = run_color_hex(run)

            if sz is not None and sz < 14.0:
                below_14.append((s_idx, sh, sz, text[:40]))
            if sz is not None and sz < 18.0:
                below_18.append((s_idx, sh, sz, text[:40]))

            if fg and bg:
                ratio = contrast(fg, bg)
                # AA-large threshold; if size < 18 OR size < 14 bold, body threshold applies.
                large = sz is not None and (sz >= 18.0 or (sz >= 14.0 and run.font.bold))
                threshold = 3.0 if large else 4.5
                if ratio < threshold:
                    contrast_violations.append(
                        (s_idx, sh, fg, bg, ratio, threshold, sz, text[:30])
                    )

    print(f'\nAudit of: {os.path.basename(pptx_path)}')
    print(f'Total text runs with content: {total_runs}')
    print(f'Runs below 14 pt: {len(below_14)}')
    print(f'Runs below 18 pt: {len(below_18)}')
    print(f'Contrast violations (fg/bg pair below threshold): {len(contrast_violations)}\n')

    if below_14:
        print('--- runs below 14pt ---')
        for s_idx, sh, sz, t in below_14:
            print(f'  slide {s_idx:>2}  {sh:<24} {sz:>5.1f}pt  "{t}"')
        print()
    if contrast_violations:
        print('--- contrast violations ---')
        for s_idx, sh, fg, bg, ratio, thr, sz, t in contrast_violations:
            print(f'  slide {s_idx:>2}  fg=#{fg} bg=#{bg}  {ratio:>4.2f}:1 (need {thr:.1f})  '
                  f'sz={sz}  "{t}"')
        print()

    fail = bool(below_14) or bool(contrast_violations)
    if fail:
        print('AUDIT FAILED — see violations above')
        return 1
    print('AUDIT PASSED — no runs below 14pt and no explicit contrast violations')
    return 0


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python scripts/audit-pptx-accessibility.py <pptx-path>')
        sys.exit(2)
    sys.exit(audit(sys.argv[1]))
