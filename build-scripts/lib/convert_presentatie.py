#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Convert 'presentatie.pptx' to a self-contained sibling HTML page.

Usage: python convert_presentatie.py [paragraph_folder_path]
       python convert_presentatie.py --all   (processes all paragraphs)

L1.5D Phase D2 — PPTX as web. Renders slide content as semantic HTML
(one <section class="slide"> per slide) with extracted text, embedded
images, and structured Vraag/Uitleg/Pitfall/Overgang speaker notes in a
collapsible <details> block per slide.

T1 (this commit) ships the converter skeleton + notes parser + minimal
single-column HTML so the deploy + landing-page pipeline is wired
end-to-end. T2 will replace `render_slide_body` with responsive
re-layout from pptx XML (kicker/title/body inference, S5 4-card cluster
detect, S6 pseudo-table detect, dark/light master classes, sticky
sidebar nav, IntersectionObserver active state). T2 changes do not
require schema changes — the slide-dict shape produced here is the
contract.
"""
import sys, io, os, glob, html as html_mod, re, hashlib
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def esc(text):
    return html_mod.escape(text if text is not None else '')


def find_paragraph_info(folder_path):
    """Extract paragraph number and name from folder path."""
    basename = os.path.basename(folder_path.rstrip('/\\'))
    parts = basename.split(' ', 1)
    number = parts[0]
    if ' - ' in basename:
        name = basename.split(' - ', 1)[1]
    else:
        name = parts[1] if len(parts) > 1 else basename
    return number, name


def get_cnvpr_descr(shape):
    """Return the descr attribute on the picture's cNvPr element, or ''."""
    try:
        xs = shape._element.xpath('.//*[local-name()="cNvPr"]/@descr')
        if xs:
            return xs[0]
    except Exception:
        pass
    return ''


def get_slide_master_name(slide):
    """Return the layout name for a slide, or '' if not retrievable.

    pptxgenjs serializes `addSlide({masterName:'DARK_HERO'})` as the
    slide LAYOUT name, not the slide-master name (the master itself is
    unnamed in pptxgenjs-built decks). The convention used in this
    project is 'LIGHT_ED' for light editorial slides and 'DARK_HERO'
    for dark hero/title/summary slides.
    """
    try:
        return slide.slide_layout.name or ''
    except Exception:
        return ''


def master_theme_class(master_name):
    """Map master name to 'slide--light' or 'slide--dark' CSS hook.

    Anything containing 'DARK' is dark; everything else (including '' and
    legacy decks) is light. Conservative default so a deck without master
    name doesn't accidentally render dark.
    """
    if 'DARK' in (master_name or '').upper():
        return 'slide--dark'
    return 'slide--light'


# ─────────────────────────────────────────────────────────────────────────────
# Notes parsing — Vraag / Uitleg / Pitfall / Overgang
# ─────────────────────────────────────────────────────────────────────────────

NOTE_KEYS = ('Vraag', 'Uitleg', 'Pitfall', 'Overgang')
_NAV_TITLE_RE = re.compile(r'^\s*NavTitle\s*:\s*(.+?)\s*$', re.MULTILINE | re.IGNORECASE)
# Matches a key prefix at the start of any line. Supports "Vraag:" with
# or without trailing whitespace. The first capture group is the key,
# the rest of the line (and continuation lines until the next key) is
# the body.
_NOTE_KEY_RE = re.compile(r'^(Vraag|Uitleg|Pitfall|Overgang)\s*:\s*(.*)$', re.MULTILINE)


def parse_notes(text):
    """Split speaker-notes text into structured dict and raw fallback.

    Returns {'structured': [(key, body), ...], 'raw': raw_text}.
    `structured` preserves source order (typically Vraag → Uitleg →
    Pitfall → Overgang but the parser tolerates any order or missing
    keys). When no canonical key is found `structured` is empty and the
    caller falls back to `raw`.

    Body text retains internal newlines; the renderer normalizes them.
    """
    if not text or not text.strip():
        return {'structured': [], 'raw': '', 'nav_title': ''}

    nav_title = ''

    def strip_nav_title(match):
        nonlocal nav_title
        if not nav_title:
            nav_title = match.group(1).strip()
        return ''

    text = _NAV_TITLE_RE.sub(strip_nav_title, text).strip()
    if not text:
        return {'structured': [], 'raw': '', 'nav_title': nav_title}

    matches = list(_NOTE_KEY_RE.finditer(text))
    if not matches:
        return {'structured': [], 'raw': text.strip(), 'nav_title': nav_title}

    structured = []
    for i, m in enumerate(matches):
        key = m.group(1)
        body_start = m.end()
        body_end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        # Re-include the first-line tail (group 2) plus continuation lines.
        first_line = m.group(2)
        rest = text[body_start:body_end]
        body = (first_line + '\n' + rest).strip('\n').rstrip()
        structured.append((key, body))
    return {'structured': structured, 'raw': text.strip(), 'nav_title': nav_title}


def render_notes(parsed):
    """Render parsed notes as a <details><dl>...</dl></details> block."""
    if not parsed['structured'] and not parsed['raw']:
        return ''
    inner_parts = []
    if parsed['structured']:
        inner_parts.append('      <dl class="slide-notes-list">')
        for key, body in parsed['structured']:
            inner_parts.append(f'        <dt class="slide-notes-key">{esc(key)}</dt>')
            inner_parts.append(
                '        <dd class="slide-notes-body">'
                + _render_notes_body(body)
                + '</dd>'
            )
        inner_parts.append('      </dl>')
    else:
        inner_parts.append(
            f'      <p class="slide-notes-fallback">{esc(parsed["raw"])}</p>'
        )
    return (
        '    <details class="slide-notes">\n'
        '      <summary>Sprekersnotities (Vraag · Uitleg · Pitfall · Overgang)</summary>\n'
        + '\n'.join(inner_parts) + '\n'
        '    </details>\n'
    )


def _render_notes_body(body):
    """Render a single notes body, preserving paragraph and indent bullets.

    Lines starting with '•' or '·' (after optional indent) become bullet
    list items; consecutive paragraphs become separate <p> tags.
    Newlines preserved as <br> within a paragraph.
    """
    if not body:
        return ''
    lines = body.split('\n')
    out_parts = []
    para_buf = []
    bullet_buf = []

    def flush_para():
        if para_buf:
            out_parts.append('<p>' + '<br>'.join(esc(l.strip()) for l in para_buf) + '</p>')
            para_buf.clear()

    def flush_bullets():
        if bullet_buf:
            out_parts.append(
                '<ul>' + ''.join(f'<li>{esc(l)}</li>' for l in bullet_buf) + '</ul>'
            )
            bullet_buf.clear()

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush_para()
            flush_bullets()
            continue
        if stripped[0] in ('•', '·'):
            flush_para()
            bullet_buf.append(stripped.lstrip('•·').strip())
        else:
            flush_bullets()
            para_buf.append(line)
    flush_para()
    flush_bullets()
    return ''.join(out_parts)


# ─────────────────────────────────────────────────────────────────────────────
# Slide walk + image extraction
# ─────────────────────────────────────────────────────────────────────────────

def extract_slide(slide, slide_idx, assets_dir, image_cache, asset_prefix='_assets'):
    """Walk one slide; return a slide dict.

    Slide dict shape (the contract T2 will refine the renderer against):
      {
        'idx': 1-based slide number,
        'master_name': raw master name string,
        'theme_class': 'slide--light' or 'slide--dark',
        'text_frames': [{'text': str, 'top_emu': int, 'left_emu': int,
                          'height_emu': int, 'width_emu': int,
                          'max_font_pt': int_or_None,
                          'all_caps': bool, 'is_decorative_empty': bool}],
        'images': [{'src': asset-relative path, 'alt': str}],
        'notes': parsed-notes dict, including optional nav_title metadata,
      }

    image_cache is a dict {sha1: (relative_src_path, alt)} so the same image
    embedded on multiple slides (e.g. we_1 on slides 6 and 7) shares one
    extracted file rather than duplicating.
    """
    text_frames = []
    images = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                blob = img.blob
                ext = img.ext or 'png'
                h = hashlib.sha1(blob).hexdigest()[:12]
                if h in image_cache:
                    src, alt = image_cache[h]
                else:
                    fname = f'presentatie-slide{slide_idx}-img-{h}.{ext}'
                    out_path = os.path.join(assets_dir, fname)
                    if not os.path.exists(out_path):
                        os.makedirs(assets_dir, exist_ok=True)
                        with open(out_path, 'wb') as f:
                            f.write(blob)
                    src = f'{asset_prefix}/{fname}'
                    alt = get_cnvpr_descr(shape) or ''
                    image_cache[h] = (src, alt)
                images.append({'src': src, 'alt': alt})
            except Exception as e:
                # Surface but don't fail the conversion — slide stays renderable.
                sys.stderr.write(
                    f'WARNING: slide {slide_idx}: image extract failed: {e}\n'
                )
            continue

        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = '\n'.join(p.text for p in tf.paragraphs).strip('\n')
        # Find largest font size used in this frame (in pt). python-pptx
        # returns None when font is inherited; treat as 0 so explicit sizes
        # dominate the title-inference heuristic in T2.
        max_pt = 0
        for p in tf.paragraphs:
            for r in p.runs:
                sz = r.font.size
                if sz is None:
                    continue
                pt = sz.pt if hasattr(sz, 'pt') else int(sz) / 12700
                if pt > max_pt:
                    max_pt = pt
        text_frames.append({
            'text': text,
            'top_emu':    shape.top if shape.top is not None else 0,
            'left_emu':   shape.left if shape.left is not None else 0,
            'height_emu': shape.height if shape.height is not None else 0,
            'width_emu':  shape.width if shape.width is not None else 0,
            'max_font_pt': max_pt or None,
            'all_caps': bool(text) and text.upper() == text and any(c.isalpha() for c in text),
            'is_decorative_empty': not text.strip(),
        })

    # Speaker notes
    notes_text = ''
    if slide.has_notes_slide:
        try:
            notes_text = slide.notes_slide.notes_text_frame.text or ''
        except Exception:
            notes_text = ''
    parsed_notes = parse_notes(notes_text)

    master_name = get_slide_master_name(slide)
    return {
        'idx': slide_idx,
        'master_name': master_name,
        'theme_class': master_theme_class(master_name),
        'text_frames': text_frames,
        'images': images,
        'notes': parsed_notes,
    }


# ─────────────────────────────────────────────────────────────────────────────
# T2 layout heuristics — kicker/title detection, card-cluster, pseudo-table
# ─────────────────────────────────────────────────────────────────────────────

# EMU constants. 1 inch = 914400 EMU. The thresholds below assume the
# 10×5.625 inch CUSTOM_16x9 layout that the §1.1.1 builder uses; they
# are wide enough to tolerate the pptxgenjs editorialSlide helper's
# fixed top-offsets for kicker (~0.3") / title (~0.6").
EMU_PER_INCH = 914400
TOP_BAND_INCH = 1.7        # frames above this are kicker/title candidates
KICKER_TOP_INCH = 0.65     # frames above this are kicker candidates
ROW_TOLERANCE_PCT = 0.55   # two frames share a row when y-midpoints differ
                           # less than this fraction of the shorter height


def detect_kicker_and_title(frames):
    """Pick kicker + title frames out of the slide's text frames.

    Returns (kicker_frame_or_None, title_frame_or_None, remaining_frames).

    Rules:
    - kicker: all-caps, short (<80 chars), top above KICKER_TOP_INCH.
    - title: tallest font in the top band (above TOP_BAND_INCH);
      tie-break by earliest top, then leftmost left. The "tallest font"
      restriction stops dramatic-display numbers in the slide body (e.g.
      slide 7's 72pt €3.500) from stealing the title slot.
    """
    kicker = None
    kicker_top_emu = KICKER_TOP_INCH * EMU_PER_INCH
    top_band_emu = TOP_BAND_INCH * EMU_PER_INCH

    for f in frames:
        if f['all_caps'] and len(f['text']) < 80 and f['top_emu'] < kicker_top_emu:
            kicker = f
            break

    title_candidates = [
        f for f in frames
        if f is not kicker and f['top_emu'] < top_band_emu
    ]
    title = None
    if title_candidates:
        title_candidates.sort(
            key=lambda f: (-(f['max_font_pt'] or 0), f['top_emu'], f['left_emu']),
        )
        title = title_candidates[0]

    remaining = [f for f in frames if f is not kicker and f is not title]
    return kicker, title, remaining


def group_rows(frames, tolerance_pct=ROW_TOLERANCE_PCT):
    """Group frames into horizontal rows by y-midpoint proximity.

    Returns a list of rows; each row is a list of frames sorted by left.
    Two frames are in the same row when their y-midpoints differ by less
    than tolerance_pct * min(height_a, height_b).
    """
    if not frames:
        return []
    sorted_frames = sorted(frames, key=lambda f: f['top_emu'])
    rows = []
    current = [sorted_frames[0]]
    for f in sorted_frames[1:]:
        prev = current[-1]
        def mid(g):
            return g['top_emu'] + g['height_emu'] / 2
        h_min = max(1, min(prev['height_emu'], f['height_emu']))
        if abs(mid(f) - mid(prev)) < tolerance_pct * h_min:
            current.append(f)
        else:
            current.sort(key=lambda g: g['left_emu'])
            rows.append(current)
            current = [f]
    current.sort(key=lambda g: g['left_emu'])
    rows.append(current)
    return rows


_CARD_NUM_RE = re.compile(r'^\s*0?[0-9]{1,2}\s*$')
_OPTION_LABEL_RE = re.compile(r'^\s*OPTIE\s+([A-Z])\s*$', re.IGNORECASE)

# A frame counts as "in the same card column" as the card-num if its
# left edge sits within this many inches to the right of the card-num.
# §1.1.1 slide 5 places cards at x=0.5"-5.5" (card num at 0.5",
# heading/subtitle at 1.65") and the Figure-3 caption at x=5.65" — well
# beyond the card column. 4 inches is generous enough to cover any
# in-card text frame but excludes side-by-side captions.
CARD_COLUMN_MAX_INCHES = 4.0
OPTION_CENTER_TOLERANCE_INCHES = 1.35
OPTION_VERTICAL_WINDOW_INCHES = 2.45


def _frame_center_x(frame):
    return frame['left_emu'] + frame['width_emu'] / 2


def detect_card_stack(frames):
    """Detect a vertically-stacked card cluster (e.g. slide 5's 4 steps).

    Returns (cards, remaining_frames). cards is a list of
    {'num': str, 'heading': str, 'subtitle': str} dicts.

    Heuristic: rows of ≥2 frames where the leftmost frame is a short
    numeric label (digit or 0-prefixed). Frames in those rows that sit
    too far to the right of the card-num (CARD_COLUMN_MAX_INCHES) are
    NOT absorbed into the card — they stay in the remainder so a
    side-by-side caption / image label doesn't pollute the last card's
    subtitle. At least 3 such consecutive rows required to count as a
    card stack. Returns ([], frames) if no match, so the caller falls
    through to pseudo-table / paragraphs.
    """
    rows = group_rows(frames)
    card_rows_filtered = []  # rows trimmed to card-column frames
    seen_cards = False
    matched_frames = set()
    max_dx_emu = int(CARD_COLUMN_MAX_INCHES * EMU_PER_INCH)

    for r in rows:
        if len(r) < 2:
            continue
        first = r[0]['text'].strip()
        if not _CARD_NUM_RE.match(first):
            continue
        card_left = r[0]['left_emu']
        # Keep only frames within the card's x-column.
        in_card = [f for f in r if (f['left_emu'] - card_left) <= max_dx_emu]
        if len(in_card) < 2:
            # The card-num is alone in its column — not a real card row.
            continue
        card_rows_filtered.append(in_card)
        seen_cards = True
        for f in in_card:
            matched_frames.add(id(f))

    if len(card_rows_filtered) < 3:
        return [], frames  # not enough cards — leave frames untouched

    cards = []
    for r in card_rows_filtered:
        num = r[0]['text'].strip()
        heading = r[1]['text'].strip() if len(r) > 1 else ''
        subtitle = ' '.join(f['text'].strip() for f in r[2:]) if len(r) > 2 else ''
        cards.append({'num': num, 'heading': heading, 'subtitle': subtitle})

    remaining = [f for f in frames if id(f) not in matched_frames]
    return cards, remaining


def detect_option_grid(frames):
    """Detect side-by-side option cards (e.g. slide 2: Optie A/B).

    PowerPoint stores each card fragment as a standalone text frame. This
    detector reconstructs the A/B pairing so the HTML preserves the comparison
    semantics instead of emitting a loose paragraph sequence.
    """
    labels = []
    for f in frames:
        m = _OPTION_LABEL_RE.match(f['text'].strip())
        if m:
            labels.append((m.group(1).upper(), f))
    if len(labels) < 2:
        return [], frames

    labels.sort(key=lambda item: item[1]['left_emu'])
    max_center_dx = int(OPTION_CENTER_TOLERANCE_INCHES * EMU_PER_INCH)
    max_vertical = int(OPTION_VERTICAL_WINDOW_INCHES * EMU_PER_INCH)
    matched = set()
    options = []

    for letter, label in labels:
        label_center = _frame_center_x(label)
        group = [
            f for f in frames
            if abs(_frame_center_x(f) - label_center) <= max_center_dx
            and f['top_emu'] >= label['top_emu'] - int(0.15 * EMU_PER_INCH)
            and f['top_emu'] <= label['top_emu'] + max_vertical
        ]
        group.sort(key=lambda f: (f['top_emu'], f['left_emu']))
        for f in group:
            matched.add(id(f))

        title = ''
        value = ''
        note_parts = []
        for f in group:
            text = f['text'].strip()
            if not text or id(f) == id(label):
                continue
            if not title:
                title = text
            elif not value and re.search(r'€|\d', text) and len(text) <= 24:
                value = text
            else:
                note_parts.append(text)

        options.append({
            'letter': letter,
            'label': f'Optie {letter}',
            'title': title,
            'value': value,
            'note': ' '.join(note_parts),
            'top_emu': label['top_emu'],
        })

    if sum(1 for o in options if o['title']) < 2:
        return [], frames

    remaining = [f for f in frames if id(f) not in matched]
    return options, remaining


def detect_pseudotable(frames):
    """Detect a 3+ column pseudo-table (e.g. slide 6's tarwe/maïs grid).

    Returns ({'header': [str], 'rows': [[str, ...]]} or None,
             remaining_frames).

    Heuristic: ≥3 rows that all have the same number of frames (≥3
    columns). Header detection: first row is treated as header iff its
    cells are all-caps OR all text frames in it share a top-band that
    sits clearly above the data rows (≥0.2" gap).
    """
    rows = group_rows(frames)
    # Find the longest run of consecutive rows with same column-count ≥3.
    best_run = []
    current_run = []
    for r in rows:
        if not current_run or len(r) == len(current_run[0]) and len(r) >= 3:
            current_run.append(r)
        else:
            if len(current_run) >= 3:
                if len(current_run) > len(best_run):
                    best_run = current_run
            current_run = [r] if len(r) >= 3 else []
    if len(current_run) >= 3 and len(current_run) > len(best_run):
        best_run = current_run

    if len(best_run) < 3:
        return None, frames

    header = None
    data_rows = best_run
    # Detect header row by EITHER majority-all-caps in the first row, OR
    # a clearly shorter row-height (column-header bands are typically
    # ~0.4" tall vs ~0.6" for data rows in the pptx authoring template).
    # The strict-all-caps rule misses cells like "TOTAAL (10 ha)" where
    # a tiny lowercase segment slips in.
    first = best_run[0]
    if len(best_run) > 1:
        caps_majority = sum(1 for f in first if f['all_caps']) * 2 >= len(first)
        first_h = sum(f['height_emu'] for f in first) / max(1, len(first))
        data_h  = sum(f['height_emu'] for r in best_run[1:] for f in r) / max(
            1, sum(len(r) for r in best_run[1:])
        )
        shorter = first_h < 0.85 * data_h  # header band is visibly thinner
        if caps_majority or shorter:
            header = [f['text'].strip() for f in first]
            data_rows = best_run[1:]

    table = {
        'header': header,
        'rows': [[f['text'].strip() for f in r] for r in data_rows],
    }

    matched = set()
    for r in best_run:
        for f in r:
            matched.add(id(f))
    remaining = [f for f in frames if id(f) not in matched]
    return table, remaining


def render_card_stack(cards):
    parts = ['    <div class="slide-cards">']
    for c in cards:
        parts.append(
            '      <div class="slide-card">\n'
            f'        <span class="card-num">{esc(c["num"])}</span>\n'
            '        <div class="card-text">\n'
            + (f'          <h3 class="card-heading">{esc(c["heading"])}</h3>\n'
               if c['heading'] else '')
            + (f'          <p class="card-subtitle">{esc(c["subtitle"])}</p>\n'
               if c['subtitle'] else '')
            + '        </div>\n'
            '      </div>'
        )
    parts.append('    </div>')
    return '\n'.join(parts)


def render_option_grid(options):
    parts = ['    <div class="slide-option-grid" role="list">']
    for opt in options:
        parts.append(
            '      <article class="slide-option-card" role="listitem">\n'
            f'        <p class="option-label">{esc(opt["label"])}</p>\n'
            + (f'        <h3 class="option-title">{esc(opt["title"])}</h3>\n'
               if opt['title'] else '')
            + (f'        <p class="option-value">{esc(opt["value"])}</p>\n'
               if opt['value'] else '')
            + (f'        <p class="option-note">{esc(opt["note"])}</p>\n'
               if opt['note'] else '')
            + '      </article>'
        )
    parts.append('    </div>')
    return '\n'.join(parts)


def render_pseudotable(table):
    parts = ['    <table class="slide-pseudotable">']
    if table['header']:
        parts.append('      <thead>')
        parts.append(
            '        <tr>'
            + ''.join(f'<th>{esc(c)}</th>' for c in table['header'])
            + '</tr>'
        )
        parts.append('      </thead>')
    parts.append('      <tbody>')
    for row in table['rows']:
        parts.append(
            '        <tr>'
            + ''.join(f'<td>{esc(c)}</td>' for c in row)
            + '</tr>'
        )
    parts.append('      </tbody>')
    parts.append('    </table>')
    return '\n'.join(parts)


def render_text_frames(frames):
    if not frames:
        return ''
    frames = sorted(frames, key=lambda f: (f['top_emu'], f['left_emu']))
    parts = ['    <div class="slide-body">']
    for f in frames:
        for line in (ln for ln in f['text'].split('\n') if ln.strip()):
            parts.append(f'      <p class="slide-text">{esc(line.strip())}</p>')
    parts.append('    </div>')
    return '\n'.join(parts)


def render_slide_body(slide):
    """Render the inner HTML of one <section class="slide">.

    T2 implementation: kicker + title pulled out by position/font; the
    remainder is passed through `detect_card_stack` and
    `detect_pseudotable` so structured blocks (S5's 4-step process,
    S6's tarwe/maïs grid) get semantic HTML. Whatever the detectors
    don't claim becomes flat <p class="slide-text"> paragraphs.
    """
    frames = [f for f in slide['text_frames'] if not f['is_decorative_empty']]

    kicker_frame, title_frame, body_frames = detect_kicker_and_title(frames)

    # Try option cards/card-stack first (more specific). Then pseudo-table on
    # the remainder.
    options, body_frames = detect_option_grid(body_frames)
    cards, body_frames = detect_card_stack(body_frames)
    table, body_frames = detect_pseudotable(body_frames)

    # Whatever survived: render in (top, left) reading order.
    body_frames.sort(key=lambda f: (f['top_emu'], f['left_emu']))

    parts = []
    if kicker_frame:
        parts.append(f'    <p class="slide-kicker">{esc(kicker_frame["text"])}</p>')
    if title_frame:
        parts.append(f'    <h2 class="slide-title">{esc(title_frame["text"])}</h2>')

    if options:
        option_top = min(o['top_emu'] for o in options)
        pre_option = [f for f in body_frames if f['top_emu'] < option_top]
        body_frames = [f for f in body_frames if f['top_emu'] >= option_top]
        pre_html = render_text_frames(pre_option)
        if pre_html:
            parts.append(pre_html)
        parts.append(render_option_grid(options))
    if cards:
        parts.append(render_card_stack(cards))
    if table:
        parts.append(render_pseudotable(table))

    body_html = render_text_frames(body_frames)
    if body_html:
        parts.append(body_html)

    for img in slide['images']:
        alt = img['alt'] or ''
        parts.append(
            f'    <figure class="slide-figure">'
            f'<img src="{esc(img["src"])}" alt="{esc(alt)}">'
            f'</figure>'
        )
    notes_html = render_notes(slide['notes'])
    if notes_html:
        parts.append(notes_html.rstrip('\n'))
    return '\n'.join(parts) + '\n'


def render_html(slides, para_number, para_name, shared_prefix='../../shared'):
    docx_filename = f'{para_number} {para_name} – presentatie.pptx'

    hero = (
        '  <header class="hero">\n'
        '    <div class="hero-inner">\n'
        '      <a class="back-link" href="../index.html"><svg viewBox="0 0 24 24"><polyline points="15 18 9 12 15 6"/></svg> Terug naar overzicht</a>\n'
        f'      <span class="hero-badge">{esc(para_number)} &middot; Presentatie</span>\n'
        f'      <h1>{esc(para_name)} &mdash; Presentatie</h1>\n'
        f'      <p class="hero-sub">{len(slides)} dia\'s &middot; les-presentatie met sprekersnotities</p>\n'
        '      <div class="hero-actions">\n'
        f'        <a class="hero-action" href="{esc(docx_filename)}" download>Download als PowerPoint</a>\n'
        '        <button type="button" class="hero-action notes-mode-toggle" data-notes-toggle aria-pressed="false">Toon alle sprekersnotities</button>\n'
        '      </div>\n'
        '    </div>\n'
        '  </header>\n'
    )

    sidebar_items = []
    for s in slides:
        # Prefer explicit slide-nav metadata from speaker notes. Fall back to
        # the legacy title-like heuristic for older decks.
        label = (s.get('notes') or {}).get('nav_title') or ''
        if not label:
            for f in s['text_frames']:
                if f['is_decorative_empty']:
                    continue
                label = f['text']
                break
            title_candidates = [
                f for f in s['text_frames']
                if not f['is_decorative_empty'] and not f['all_caps']
            ]
            if title_candidates:
                title_candidates.sort(
                    key=lambda f: (-(f['max_font_pt'] or 0), f['top_emu']),
                )
                label = title_candidates[0]['text']
        label_short = (label or 'Dia').split('\n')[0]
        if len(label_short) > 32:
            label_short = label_short[:31] + '…'
        sidebar_items.append(
            f'        <a class="nav-item" href="#slide-{s["idx"]}" data-slide-link="{s["idx"]}">'
            f'<span class="nav-num">{s["idx"]:02d}</span>'
            f'<span class="nav-title">{esc(label_short)}</span></a>'
        )

    sidebar = (
        '  <nav class="sidebar" id="sidebar">\n'
        '    <div class="sidebar-header">\n'
        f'      <h2>{esc(para_number)} {esc(para_name)}</h2>\n'
        '      <p>Presentatie</p>\n'
        '    </div>\n'
        + '\n'.join(sidebar_items) + '\n'
        '  </nav>\n'
    )

    section_blocks = []
    for s in slides:
        body = render_slide_body(s)
        section_blocks.append(
            f'      <section class="slide {s["theme_class"]}" id="slide-{s["idx"]}" data-slide="{s["idx"]}" hidden>\n'
            f'{body}'
            f'      </section>'
        )

    nav_bar = (
        '    <div class="slide-controls" role="navigation" aria-label="Dia-navigatie">\n'
        '      <button type="button" class="slide-prev" aria-label="Vorige dia" data-slide-prev>&larr; Vorige</button>\n'
        f'      <span class="slide-counter"><span data-slide-current>1</span> / <span class="slide-total">{len(slides)}</span></span>\n'
        '      <button type="button" class="slide-next" aria-label="Volgende dia" data-slide-next>Volgende &rarr;</button>\n'
        '    </div>\n'
    )

    main_grid = (
        '    <main class="presentatie-grid">\n'
        + nav_bar
        + '      <div class="presentatie-deck" data-slide-deck>\n'
        + '\n'.join(section_blocks) + '\n'
        + '      </div>\n'
        + nav_bar
        + '    </main>\n'
    )

    boot_script = (
        '<script>\n'
        '(function(){try{var m=localStorage.getItem("quizMode");\n'
        'if(!m && window.matchMedia && window.matchMedia("(prefers-color-scheme: dark)").matches){m="dark";}\n'
        'if(m==="dark"){document.documentElement.setAttribute("data-theme","dark");}\n'
        '}catch(e){}})();\n'
        '</script>'
    )

    slide_nav_script = '''<script>
(function(){
  var deck = document.querySelector('[data-slide-deck]');
  if (!deck) return;
  var slides = Array.prototype.slice.call(deck.querySelectorAll('.slide'));
  if (!slides.length) return;
  var sidebarLinks = document.querySelectorAll('[data-slide-link]');
  var counters = document.querySelectorAll('[data-slide-current]');
  var prevButtons = document.querySelectorAll('[data-slide-prev]');
  var nextButtons = document.querySelectorAll('[data-slide-next]');
  var notesToggle = document.querySelector('[data-notes-toggle]');
  var current = 0;

  function show(idx){
    if (idx < 0) idx = 0;
    if (idx >= slides.length) idx = slides.length - 1;
    slides.forEach(function(s, i){
      if (i === idx) { s.hidden = false; s.classList.add('is-active'); }
      else { s.hidden = true; s.classList.remove('is-active'); }
    });
    sidebarLinks.forEach(function(a){
      var n = parseInt(a.getAttribute('data-slide-link'), 10);
      if (n === idx + 1) a.classList.add('is-active');
      else a.classList.remove('is-active');
    });
    counters.forEach(function(el){ el.textContent = String(idx + 1); });
    prevButtons.forEach(function(b){ b.disabled = (idx === 0); });
    nextButtons.forEach(function(b){ b.disabled = (idx === slides.length - 1); });
    current = idx;
    var slideId = slides[idx].id;
    if (slideId && history.replaceState) {
      history.replaceState(null, '', '#' + slideId);
    }
  }

  function fromHash(){
    var m = (location.hash || '').match(/^#slide-(\\d+)$/);
    if (m) {
      var n = parseInt(m[1], 10) - 1;
      if (n >= 0 && n < slides.length) return n;
    }
    return 0;
  }

  prevButtons.forEach(function(b){ b.addEventListener('click', function(){ show(current - 1); }); });
  nextButtons.forEach(function(b){ b.addEventListener('click', function(){ show(current + 1); }); });

  sidebarLinks.forEach(function(a){
    a.addEventListener('click', function(e){
      var n = parseInt(a.getAttribute('data-slide-link'), 10);
      if (!isNaN(n)) {
        e.preventDefault();
        show(n - 1);
      }
    });
  });

  if (notesToggle) {
    notesToggle.addEventListener('click', function(){
      var details = Array.prototype.slice.call(document.querySelectorAll('details.slide-notes'));
      var shouldOpen = details.some(function(d){ return !d.open; });
      details.forEach(function(d){ d.open = shouldOpen; });
      notesToggle.setAttribute('aria-pressed', shouldOpen ? 'true' : 'false');
      notesToggle.textContent = shouldOpen ? 'Verberg alle sprekersnotities' : 'Toon alle sprekersnotities';
    });
  }

  document.addEventListener('keydown', function(e){
    if (e.target && /^(INPUT|TEXTAREA|SELECT)$/.test(e.target.tagName)) return;
    if (e.key === 'ArrowRight' || e.key === 'PageDown') { e.preventDefault(); show(current + 1); }
    else if (e.key === 'ArrowLeft' || e.key === 'PageUp') { e.preventDefault(); show(current - 1); }
    else if (e.key === 'Home') { e.preventDefault(); show(0); }
    else if (e.key === 'End') { e.preventDefault(); show(slides.length - 1); }
  });

  window.addEventListener('hashchange', function(){ show(fromHash()); });

  show(fromHash());
})();
</script>'''

    return f'''<!doctype html>
<html lang="nl" data-theme="light">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{esc(para_name)} &mdash; Presentatie</title>
  <link rel="stylesheet" href="{shared_prefix}/voorkennis.css">
  {boot_script}
</head>
<body data-layout="presentatie-v1" data-accent-domain="economisch">
<button class="sidebar-toggle" id="sidebarToggle" aria-label="Menu">&#9776;</button>
<div class="sidebarOverlay" id="sidebarOverlay"></div>
<div class="page-layout">
{sidebar}  <div class="content">
{hero}{main_grid}  </div>
</div>
<script src="{shared_prefix}/voorkennis.js"></script>
{slide_nav_script}
</body>
</html>
'''


# ─────────────────────────────────────────────────────────────────────────────
# Per-paragraph processing + CLI
# ─────────────────────────────────────────────────────────────────────────────

def process_paragraph(para_folder):
    """Process a single paragraph folder.

    Returns 'ok' | 'skip' | 'error'. 'skip' means no presentatie.pptx
    found (acceptable — paragraph doesn't carry this surface yet).
    'error' means parse/IO failure — deploy.js must surface as exit 1.
    """
    para_number, para_name = find_paragraph_info(para_folder)
    exact_name = f'{para_number} {para_name} – presentatie.pptx'
    exact_path = os.path.join(para_folder, exact_name)
    if os.path.exists(exact_path):
        found = [exact_path]
    else:
        pattern = os.path.join(para_folder, '*presentatie*.pptx')
        found = [
            f for f in glob.glob(pattern)
            if not os.path.basename(f).startswith('~$')
            and 'prototype' not in os.path.basename(f).lower()
        ]
        found.sort()
    if not found:
        print(f'  SKIP {para_number}: no presentatie.pptx found')
        return 'skip'
    pptx_path = found[0]

    assets_dir = os.path.join(para_folder, '_assets')

    try:
        pres = Presentation(pptx_path)
    except Exception as e:
        print(f'  ERROR {para_number}: cannot open pptx: {e}')
        return 'error'

    image_cache = {}
    slides = []
    try:
        for i, slide in enumerate(pres.slides, 1):
            slides.append(extract_slide(slide, i, assets_dir, image_cache))
    except Exception as e:
        print(f'  ERROR {para_number}: slide parse failed: {e}')
        return 'error'

    if not slides:
        print(f'  SKIP {para_number}: pptx has zero slides')
        return 'skip'

    try:
        html_content = render_html(slides, para_number, para_name)
        html_path = os.path.join(
            para_folder, f'{para_number} {para_name} – presentatie.html'
        )
        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
    except Exception as e:
        print(f'  ERROR {para_number}: html render failed: {e}')
        return 'error'

    # Surface alt-text gaps as a WARN (not error) so the deploy keeps going
    # but the author sees the issue.
    empty_alt = sum(1 for s in slides for img in s['images'] if not img['alt'])
    n_images  = sum(len(s['images']) for s in slides)
    n_notes_struct = sum(1 for s in slides if s['notes']['structured'])
    if empty_alt:
        sys.stderr.write(
            f'WARNING: {para_number}: {empty_alt}/{n_images} slide images '
            'have empty alt-text (descr); add altText to addImage calls in '
            'the paragraph pptx builder.\n'
        )
    print(
        f'  OK {para_number} {para_name} ({len(slides)} slides, '
        f'{n_images} images, {n_notes_struct} structured-notes)'
    )
    return 'ok'


def main():
    args = sys.argv[1:]
    if not args:
        print('Usage: python convert_presentatie.py [path] or --all')
        sys.exit(1)

    had_error = False
    if args[0] == '--all':
        base = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        siblings_root = os.path.dirname(base)
        folders = []
        for boek in sorted(glob.glob(os.path.join(siblings_root, '4veco-lessen', 'Boek*'))):
            for chap in sorted(glob.glob(os.path.join(boek, '*'))):
                if not os.path.isdir(chap):
                    continue
                for para in sorted(glob.glob(os.path.join(chap, '*'))):
                    if os.path.isdir(para):
                        folders.append(para)
        for f in sorted(glob.glob(os.path.join(base, '3.*/3.*.*'))):
            folders.append(f)
        print(f'Found {len(folders)} candidate paragraph folders')
        success = 0
        for folder in folders:
            result = process_paragraph(folder)
            if result == 'ok':
                success += 1
            elif result == 'error':
                had_error = True
        print(f'\nDone: {success}/{len(folders)} converted')
    else:
        result = process_paragraph(args[0])
        if result == 'error':
            had_error = True

    sys.exit(1 if had_error else 0)


if __name__ == '__main__':
    main()
